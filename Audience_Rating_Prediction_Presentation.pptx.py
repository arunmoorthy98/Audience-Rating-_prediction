from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

# Create a new presentation
presentation = Presentation()

# Function to add gradient background
def add_gradient_background(slide):
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(70, 130, 180)  # Steel Blue
    fill.gradient_stops[1].color.rgb = RGBColor(255, 228, 181)  # Moccasin

# Function to add styled title
def add_title_with_shape(slide, text):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(9), Inches(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(25, 25, 112)  # Midnight Blue
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.paragraphs[0].font.size = Pt(24)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Slide 1: Title Slide with Gradient Background
slide = presentation.slides.add_slide(presentation.slide_layouts[0])
add_gradient_background(slide)

title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Movie Audience Rating Prediction"
subtitle.text = "Building a predictive model and validating accuracy"
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Slide 2: Problem Statement with Icon Bullets
slide = presentation.slides.add_slide(presentation.slide_layouts[5])  # Blank slide
add_gradient_background(slide)

add_title_with_shape(slide, "Problem Statement")
content = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
text_frame = content.text_frame
text_frame.text = (
    "- Objective: Build a model to predict 'audience_rating'.\n"
    "- Dataset: Rotten Tomatoes movies dataset with features like rating, genre, and directors."
)
for paragraph in text_frame.paragraphs:
    paragraph.font.size = Pt(18)
    paragraph.font.color.rgb = RGBColor(0, 0, 0)

# Slide 3: Data Overview with Shapes
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
add_gradient_background(slide)

add_title_with_shape(slide, "Data Overview")
content = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
text_frame = content.text_frame
text_frame.text = (
    "Dataset Details:\n"
    "- Size: Includes both numerical and categorical features.\n"
    "- Target Variable: audience_rating.\n"
    "- Key Features: rating, genre, runtime, tomatometer_rating."
)
for paragraph in text_frame.paragraphs:
    paragraph.font.size = Pt(16)
    paragraph.font.color.rgb = RGBColor(0, 51, 102)

# Slide 4: Data Preprocessing with Icons
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
add_gradient_background(slide)

add_title_with_shape(slide, "Data Preprocessing")
content = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
text_frame = content.text_frame
text_frame.text = (
    "1. Handle Missing Values:\n"
    "   - Numerical: Impute with mean.\n"
    "   - Categorical: Impute with most frequent value.\n"
    "2. Standardize Features: Use StandardScaler.\n"
    "3. Encode Categorical Data: One-Hot Encoding."
)
for paragraph in text_frame.paragraphs:
    paragraph.font.size = Pt(16)
    paragraph.font.color.rgb = RGBColor(0, 51, 51)

# Slide 5: Model Building with a Footer
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
add_gradient_background(slide)

add_title_with_shape(slide, "Model Building")
content = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
text_frame = content.text_frame
text_frame.text = (
    "- Model: Random Forest Regressor.\n"
    "- Pipeline: Integrated preprocessing and training.\n"
    "- Data Split: 80% training, 20% testing."
)
for paragraph in text_frame.paragraphs:
    paragraph.font.size = Pt(16)
    paragraph.font.color.rgb = RGBColor(102, 0, 51)

footer = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
footer.text = "Model built using Scikit-learn with integrated pipelines."
footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(105, 105, 105)
footer.text_frame.paragraphs[0].font.size = Pt(12)

# Slide 6: Model Validation
slide = presentation.slides.add_slide(presentation.slide_layouts[5])
add_gradient_background(slide)

add_title_with_shape(slide, "Model Validation")
content = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
text_frame = content.text_frame
text_frame.text = (
    "Evaluation Metrics:\n"
    "- Mean Absolute Error (MAE): 11.49\n"
    "- Mean Squared Error (MSE): 217.50\n"
    "- Root Mean Squared Error (RMSE): 14.75\n"
    "- R-squared (R²): 0.67"
)
for paragraph in text_frame.paragraphs:
    paragraph.font.size = Pt(16)
    paragraph.font.color.rgb = RGBColor(139, 69, 19)

# Save the presentation
output_path = "Enhanced_Audience_Rating_Presentation.pptx"
presentation.save(output_path)
print(f"Presentation saved as {output_path}")
